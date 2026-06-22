"""
Gera a apresentação de churn prediction em .pptx para o cliente (não técnico).
Executar: python reports/gerar_apresentacao.py   (da raiz do projeto)
Output:   reports/apresentacao_churn.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── Paleta ────────────────────────────────────────────────────────────────────
V_ESCURO  = RGBColor(0x1b, 0x43, 0x32)   # #1b4332
V_MEDIO   = RGBColor(0x40, 0x91, 0x6c)   # #40916c
V_SUAVE   = RGBColor(0x52, 0xb7, 0x88)   # #52b788
V_CLARO   = RGBColor(0xd8, 0xf3, 0xdc)   # #d8f3dc
BRANCO    = RGBColor(0xFF, 0xFF, 0xFF)
CINZA     = RGBColor(0x4a, 0x4a, 0x4a)
VERMELHO  = RGBColor(0xe6, 0x39, 0x46)

# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=V_ESCURO,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def slide_background(slide, color=BRANCO):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_bullet_box(slide, items, left, top, width, height,
                   font_size=16, color=CINZA, bullet="●  "):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

def header_bar(slide, title, subtitle=None):
    """Barra verde escura no topo com título."""
    add_rect(slide, 0, 0, 10, 1.3, V_ESCURO)
    add_textbox(slide, title, 0.3, 0.1, 9, 0.7,
                font_size=28, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.78, 9, 0.4,
                    font_size=14, color=V_CLARO, align=PP_ALIGN.LEFT)

def accent_line(slide, top=1.35):
    add_rect(slide, 0, top, 10, 0.05, V_SUAVE)

def kpi_card(slide, valor, label, left, top, width=2.1, height=1.2):
    add_rect(slide, left, top, width, height, V_CLARO)
    add_rect(slide, left, top, width, 0.04, V_SUAVE)
    add_textbox(slide, valor, left + 0.1, top + 0.1, width - 0.2, 0.6,
                font_size=26, bold=True, color=V_ESCURO, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left + 0.05, top + 0.65, width - 0.1, 0.5,
                font_size=12, color=CINZA, align=PP_ALIGN.CENTER)

# ── Apresentação ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

BLANK = prs.slide_layouts[6]   # layout em branco

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_background(s, V_ESCURO)

add_rect(s, 0, 0, 10, 5.625, V_ESCURO)
add_rect(s, 0, 3.8, 10, 1.825, V_MEDIO)
add_rect(s, 0.5, 1.8, 0.08, 1.6, V_SUAVE)

add_textbox(s, "Previsão de Churn", 0.8, 1.2, 8.5, 1.1,
            font_size=40, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)
add_textbox(s, "de Clientes", 0.8, 2.0, 8.5, 0.8,
            font_size=40, bold=True, color=V_SUAVE, align=PP_ALIGN.LEFT)
add_textbox(s, "Identificando em quem investir esforço de retenção",
            0.8, 2.9, 8.5, 0.6, font_size=17, color=V_CLARO, align=PP_ALIGN.LEFT)
add_textbox(s, f"Junho de 2026", 0.8, 4.1, 4, 0.5,
            font_size=13, color=BRANCO, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — O Desafio
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_background(s)
header_bar(s, "O Desafio", "Por que prever quem vai parar de comprar?")
accent_line(s)

add_textbox(s,
    "Clientes que somem custam muito mais do que clientes que ficam.",
    0.4, 1.5, 9.2, 0.6, font_size=18, bold=True, color=V_ESCURO)

add_bullet_box(s, [
    "Reconquistar um cliente perdido custa 5–7× mais do que reter um cliente ativo",
    "Clientes B2B com contrato têm alto valor recorrente — cada saída impacta diretamente o faturamento",
    "A saída raramente é abrupta: o cliente reduz gradualmente antes de desaparecer",
    "Identificar esse padrão cedo permite agir com antecedência — antes que a perda aconteça",
], 0.5, 2.1, 9, 2.8, font_size=15)

add_rect(s, 0, 5.2, 10, 0.425, V_CLARO)
add_textbox(s, "💡  Pergunta central:  Quais clientes têm maior risco de não comprar nos próximos 3 meses?",
            0.3, 5.22, 9.4, 0.35, font_size=13, bold=True, color=V_ESCURO)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Base Analisada
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_background(s)
header_bar(s, "Base Analisada", "O que alimenta o modelo")
accent_line(s)

kpi_card(s, "765",          "clientes com\npedidos desde 2023",   0.35, 1.55)
kpi_card(s, "535 mil",      "registros de\npedidos analisados",   2.65, 1.55)
kpi_card(s, "3+ anos",      "de histórico\ncomportamental",       4.95, 1.55)
kpi_card(s, "1.537",        "itens diferentes\nno portfólio",     7.25, 1.55)

add_textbox(s, "Tipos de dado utilizados", 0.4, 3.0, 9, 0.4,
            font_size=15, bold=True, color=V_ESCURO)
add_bullet_box(s, [
    "Frequência de compra — quantas vezes e com que regularidade cada cliente pediu",
    "Valor e portfólio — ticket médio, itens comprados, tendência de queda por produto",
    "Recência — há quantos meses o cliente não faz pedidos",
    "Inadimplência — histórico de dias em atraso no pagamento",
], 0.5, 3.4, 9.2, 2.0, font_size=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Perfil de Risco por Categoria
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_background(s)
header_bar(s, "Perfil dos Clientes", "Nem todos os clientes têm o mesmo risco")
accent_line(s)

# Tabela simples de categorias
categorias = [
    ("Premium",  "≥ 20 pedidos/mês",  "20%",  V_ESCURO),
    ("Alto",     "6–19 pedidos/mês",  "24%",  V_MEDIO),
    ("Médio",    "3–5 pedidos/mês",   "11%",  V_SUAVE),
    ("Baixo",    "1–2 pedidos/mês",   "40%",  rgb(0x95, 0xd5, 0xb2)),
]
add_textbox(s, "Categoria", 0.4, 1.5, 2, 0.35, font_size=13, bold=True, color=CINZA)
add_textbox(s, "Volume",    2.5, 1.5, 2, 0.35, font_size=13, bold=True, color=CINZA)
add_textbox(s, "Taxa de churn observada", 4.6, 1.5, 3, 0.35, font_size=13, bold=True, color=CINZA)
add_rect(s, 0.35, 1.85, 9.3, 0.03, rgb(0xcc, 0xcc, 0xcc))

for i, (cat, vol, taxa, cor) in enumerate(categorias):
    y = 2.0 + i * 0.7
    add_rect(s, 0.35, y, 0.18, 0.45, cor)
    add_textbox(s, cat,  0.65, y + 0.05, 1.8, 0.4, font_size=14, bold=True,  color=V_ESCURO)
    add_textbox(s, vol,  2.5,  y + 0.05, 2.0, 0.4, font_size=13, color=CINZA)
    add_textbox(s, taxa, 4.6,  y + 0.05, 1.0, 0.4, font_size=20, bold=True,  color=cor if cor != rgb(0x95, 0xd5, 0xb2) else V_MEDIO)
    bar_w = float(taxa.replace("%","")) / 100 * 3.5
    add_rect(s, 5.8, y + 0.12, bar_w, 0.28, cor)

add_rect(s, 0, 4.95, 10, 0.675, V_CLARO)
add_textbox(s,
    "Clientes de volume baixo são a maioria da base e apresentam a maior taxa de churn — "
    "são o principal alvo do modelo.",
    0.3, 5.0, 9.4, 0.55, font_size=13, color=V_ESCURO)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Como Funciona o Modelo
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_background(s)
header_bar(s, "Como Funciona", "O modelo aprende com o passado para antecipar o futuro")
accent_line(s)

etapas = [
    ("1", "Histórico",     "Analisamos 3 anos de comportamento de compra de cada cliente"),
    ("2", "Padrão",        "O modelo identifica o perfil dos clientes que pararam de comprar"),
    ("3", "Comparação",    "Compara cada cliente ativo com esse padrão histórico"),
    ("4", "Score de risco","Cada cliente recebe uma pontuação de 0 a 1 — quanto maior, maior o risco"),
]

for i, (num, titulo, desc) in enumerate(etapas):
    x = 0.4 + i * 2.38
    add_rect(s, x, 1.55, 2.1, 2.5, V_CLARO)
    add_rect(s, x, 1.55, 2.1, 0.55, V_ESCURO)
    add_textbox(s, num,    x + 0.05, 1.57, 0.5,  0.45, font_size=24, bold=True, color=V_SUAVE)
    add_textbox(s, titulo, x + 0.55, 1.62, 1.5,  0.4,  font_size=14, bold=True, color=BRANCO)
    add_textbox(s, desc,   x + 0.1,  2.2,  1.9,  1.7,  font_size=12, color=CINZA)

    if i < 3:
        add_textbox(s, "→", x + 2.12, 2.45, 0.25, 0.4, font_size=20, bold=True, color=V_MEDIO)

add_rect(s, 0, 4.35, 10, 0.5, V_ESCURO)
add_textbox(s,
    "Resultado: lista mensal com os clientes que mais precisam de atenção — antes que a perda aconteça.",
    0.3, 4.38, 9.4, 0.42, font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — O Que Indica Risco
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_background(s)
header_bar(s, "O Que Indica Risco de Churn", "Os sinais que o modelo aprendeu a reconhecer")
accent_line(s)

sinais = [
    ("🕐", "Tempo sem pedir",
     "O sinal mais forte. Clientes que param de comprar por períodos mais longos que seu padrão histórico entram em alerta."),
    ("📉", "Redução de itens no portfólio",
     "Antes de sumir, o cliente tende a reduzir gradualmente a variedade de produtos — sinal de migração para concorrente."),
    ("📆", "Intervalo entre compras crescendo",
     "Clientes que historicamente pediam a cada 30 dias e passam a pedir a cada 60 estão sinalizando afastamento."),
    ("⚠️", "Histórico de inadimplência",
     "Dias em atraso no pagamento correlacionam com risco de churn — clientes com atrito financeiro saem mais."),
]

for i, (icone, titulo, desc) in enumerate(sinais):
    row, col = divmod(i, 2)
    x = 0.35 + col * 4.85
    y = 1.55 + row * 1.8
    add_rect(s, x, y, 4.5, 1.55, V_CLARO)
    add_rect(s, x, y, 0.55, 1.55, V_ESCURO)
    add_textbox(s, icone,  x + 0.07, y + 0.42, 0.5,  0.5,  font_size=18, color=BRANCO)
    add_textbox(s, titulo, x + 0.65, y + 0.1,  3.75, 0.4,  font_size=14, bold=True, color=V_ESCURO)
    add_textbox(s, desc,   x + 0.65, y + 0.52, 3.75, 1.0,  font_size=11, color=CINZA)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Desempenho do Modelo
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_background(s)
header_bar(s, "Desempenho do Modelo", "Validado em dados reais que o modelo nunca viu")
accent_line(s)

add_textbox(s,
    "Testamos o modelo em 724 clientes com histórico de 2023–2025, "
    "verificando se as previsões para os 3 meses seguintes se confirmaram.",
    0.4, 1.5, 9.2, 0.65, font_size=14, color=CINZA)

kpi_card(s, "91%",   "precisão na\nordenação de risco*",    0.4,  2.35, 2.8)
kpi_card(s, "3 em 4","alertas são clientes\nque realmente saíram", 3.4, 2.35, 3.0)
kpi_card(s, "9 em 10","churners capturados\npelo modelo",   6.6,  2.35, 2.8)

add_textbox(s,
    "* O modelo ordena corretamente o risco em 91% dos pares de clientes comparados.",
    0.4, 3.8, 9, 0.35, font_size=11, color=CINZA, italic=True)

add_rect(s, 0, 4.2, 10, 0.5, V_CLARO)
add_textbox(s,
    "Comparação: um modelo aleatório acertaria apenas 50% — o modelo chega a 91%.",
    0.4, 4.25, 9.2, 0.4, font_size=13, color=V_ESCURO)

add_textbox(s,
    "Por que não 100%? Alguns churners são imprevisíveis — decisões externas ao comportamento de compra "
    "(mudança de fornecedor, fechamento do negócio). O modelo captura o que os dados conseguem explicar.",
    0.4, 4.8, 9.2, 0.7, font_size=11, color=CINZA, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Como Usar os Alertas
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_background(s)
header_bar(s, "Como Usar os Alertas", "Da lista ao contato — um processo simples")
accent_line(s)

passos = [
    ("Receber a lista",
     "Todo mês uma planilha com os clientes ordenados por score de risco é gerada automaticamente."),
    ("Priorizar pelo score",
     "Clientes com score mais alto = maior urgência. Comece pelos primeiros da lista."),
    ("Entender o motivo",
     "A planilha indica o principal sinal de alerta de cada cliente (tempo sem pedir, queda de portfólio...)."),
    ("Agir preventivamente",
     "Contato proativo: oferta comercial, visita, ou simplesmente uma ligação — antes de perder o cliente."),
]

for i, (titulo, desc) in enumerate(passos):
    y = 1.5 + i * 0.9
    add_rect(s, 0.35, y, 0.55, 0.65, V_ESCURO)
    add_textbox(s, str(i+1), 0.35, y + 0.08, 0.55, 0.5,
                font_size=20, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(s, titulo, 1.05, y + 0.05, 2.5, 0.35,
                font_size=14, bold=True, color=V_ESCURO)
    add_textbox(s, desc,   1.05, y + 0.38, 8.5, 0.45,
                font_size=12, color=CINZA)
    if i < 3:
        add_rect(s, 0.52, y + 0.65, 0.2, 0.25, V_SUAVE)

add_rect(s, 0, 5.1, 10, 0.525, V_CLARO)
add_textbox(s,
    "📄  Arquivo entregue:  clientes_alerta_modelo.xlsx  —  nome do cliente, score de risco, principal sinal de alerta, categoria",
    0.3, 5.13, 9.4, 0.38, font_size=12, color=V_ESCURO)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Próximos Passos
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_background(s)
header_bar(s, "Próximos Passos", "Para colocar o modelo em operação")
accent_line(s)

steps = [
    ("Curto prazo",  V_ESCURO, [
        "Validar a lista de alertas com a equipe comercial — o modelo acertou os clientes que vocês já suspeitavam?",
        "Definir capacidade de retenção: quantos clientes a equipe consegue contatar por mês?",
        "Confirmar custo relativo: perder um cliente vs. abordar um que não ia sair",
    ]),
    ("Médio prazo", V_MEDIO, [
        "Atualização mensal automática da lista de alertas com dados novos",
        "Registro das ações tomadas — quais clientes foram contatados e qual foi o resultado",
        "Refinamento do modelo com feedback das ações de retenção",
    ]),
]

for i, (periodo, cor, itens) in enumerate(steps):
    x = 0.35 + i * 4.85
    add_rect(s, x, 1.5, 4.5, 0.55, cor)
    add_textbox(s, periodo, x + 0.2, 1.55, 4.1, 0.45,
                font_size=16, bold=True, color=BRANCO)
    add_bullet_box(s, itens, x + 0.1, 2.15, 4.3, 3.0,
                   font_size=12, color=CINZA)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Encerramento
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_background(s, V_ESCURO)

add_rect(s, 0, 0, 10, 5.625, V_ESCURO)
add_rect(s, 0, 2.2, 10, 1.2, V_MEDIO)
add_rect(s, 3.5, 0.8, 0.08, 4.0, V_SUAVE)

add_textbox(s, "Obrigada.", 3.8, 0.6, 5.8, 1.0,
            font_size=40, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)
add_textbox(s,
    "O modelo está pronto.\nA lista de clientes em risco também.",
    3.8, 2.3, 5.8, 0.9, font_size=18, color=V_CLARO, align=PP_ALIGN.LEFT)
add_textbox(s,
    "Agora é hora de agir.",
    3.8, 3.3, 5.8, 0.6, font_size=16, color=V_SUAVE, align=PP_ALIGN.LEFT,
    italic=True)

# ── Salvar ────────────────────────────────────────────────────────────────────
from pathlib import Path
output = Path(__file__).parent / "apresentacao_churn.pptx"
prs.save(output)
print(f"OK: Apresentacao salva em: {output}")
print(f"  Slides: {len(prs.slides)}")
